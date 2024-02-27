from pptx import Presentation 
from pptx.util import Inches, Pt

root = Presentation('Ex.pptx')

for slide in root.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.italic = True

root.save('Ex_2.pptx')