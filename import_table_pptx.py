from pptx import Presentation 
from pptx.util import Inches


root = Presentation()

black_slide = root.slide_layouts[6]

slide = root.slides.add_slide(black_slide)

left = top = Inches(1)
width = height = Inches(5) # Размер таблицы

shape = slide.shapes.add_table(3, 4, left, top, width, height)

table = shape.table

cell = table.cell(0, 1)
cell.text = 'Column_1'

root.save('Ex3.pptx')