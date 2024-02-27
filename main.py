from pptx import Presentation 
from pptx.util import Inches, Pt

root = Presentation()

############### First slide

first_slide_layout = root.slide_layouts[0]

slide = root.slides.add_slide(first_slide_layout)

slide.shapes.title.text = 'Тест'

slide.placeholders[1].text = 'Подзаголовок Тест'


############### Second slide

second_slide_layout = root.slide_layouts[6]

slide_2 = root.slides.add_slide(second_slide_layout)

left = top = Inches(1) # Отступы от края листа
width = height = Inches(5) # Размер рамки

txBox = slide_2.shapes.add_textbox(left, top, width, height)

tf = txBox.text_frame

# Настройки одной строки текста
p = tf.add_paragraph()
p.text = 'This is an Example text This is an Example text'
p.font.bold = True
p.font.italic = True
p.front.name = 'Arial'
p.font.size = Pt(50)

# Добавить ещё одну строку
p = tf.add_paragraph()
p.text = 'This is an Example text 2'


root.save('Examlpe.pptx')


