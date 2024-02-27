from pptx import Presentation 
from pptx.util import Inches

img_path = 'img.PNG'

root = Presentation()

black_slide = root.slide_layouts[6]

slide = root.slides.add_slide(black_slide)

left = top = Inches(1)
width = height = Inches(5) # Размер изображения в дюймах

pic = slide.shapes.add_picture(img_path, left, top, width, height)

root.save('Ex3.pptx')